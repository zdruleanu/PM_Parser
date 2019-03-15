import pandas as pd
import os
import shutil
#import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from bokeh.plotting import figure, show
from bokeh.io import output_file
from bokeh.models.widgets import Tabs, Panel
from bokeh.io import export_png
import bokeh.palettes as palletes
from bokeh.models import Legend, LegendItem
import itertools
import yaml
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from sys import argv
from matplotlib.widgets import TextBox


class GraficCollection:
    '''
    This class is used to contain the data from one single graph(figure)
    '''
    def __init__(self, figura):
        self.figura = figura
        self.grafice = {}

    def extrage_gafice(self, excel_file, filter_params, target_columns, use_target_columns_in_legend):
        # The function returns a list of pandas series, where each series will be ploted on the same plot
        # ExcelFile is the imported DataFrame (already indexed by time)
        # FilterParams is a dictionary with keys being column names and values being a tuple, where the first item is
        #   the operator, the second being a list of values that each column should be filtered upon,
        # TargetColumns is a list of column names that we want in the graph.
        # when we plot it will always plot based on the ExcelFie indexes on x axis
        #
        # FilterParams = [['Slot', '==',[0,1,4]], ['altceva', '==',['d','d']], ['al3', '<=',[7]]]

        FilterParamNumber = filter_params.__len__()
        AllColValuesLists = [i[2] for i in filter_params]
        AllColValuesListsCombinations = list(itertools.product(*AllColValuesLists))

        while True:
            try:
                comb = list(AllColValuesListsCombinations.pop(0))
            except:
                print('Finished graphs data generation')
                break
            FilterExpression = '('
            for i in range(0,FilterParamNumber):
                print(comb[i])
                if isinstance(comb[i], str):
                    comb[i] = "'"+comb[i]+"'"
                FilterExpression += "(excel_file['" + filter_params[i][0] + "'] " + filter_params[i][1] + \
                                    " " + str(comb[i]) + ")"
                if i < FilterParamNumber-1:
                    FilterExpression += " & "
                else:
                    FilterExpression += ")"
            print("FiltereExpression is: ", FilterExpression)
            filterResult = excel_file[eval(FilterExpression)]
            PrefixNumeGrafic = ''
            for FilterParam in comb:
                PrefixNumeGrafic += str(FilterParam) + "_"
            # We consider the first filter parameter relevnat for the name of the grapsh as it will be shown in the 
            # legend. For example it will represent the Slot number
            #PrefixNumeGrafic = str(comb[0])

            for TargetColumn in target_columns:
                if not filterResult[TargetColumn].empty:
                    if use_target_columns_in_legend:
                        self.grafice[PrefixNumeGrafic + "_" + TargetColumn] = (filterResult[TargetColumn])
                    else:
                        # Filtering parameters might lead to a result where the PrefixNumeGrafic is not unique.
                        # If so, exit and print Error
                        if PrefixNumeGrafic in self.grafice.keys():
                            print("ERROR: graph names are not unique. Please set UseTargetColumnsInLegend "
                                  "to True in the config file")
                            exit()
                        else:
                            self.grafice[PrefixNumeGrafic] = (filterResult[TargetColumn])


def find_data_filenames(path_to_dir, prefix="", suffix=""):
    # returns a list of file names which include prefix and suffix.
    # the file names returned include the path
    # directories are not returned

    filenames = [path_to_dir+"/"+filename for filename in os.listdir(path_to_dir) if (filename.startswith(prefix) and
                 filename.endswith(suffix) and os.path.isfile(path_to_dir+"/"+filename))]
    return filenames


def check_if_template(extracted_for_template_check, template_dir):
    # Checks in the templateDir if there exists a file whos name begins with extracted_for_template_check.
    # Returns the template file name (with path)
    print("Checking for tempates")
    templateFileName = find_data_filenames(template_dir, prefix=extracted_for_template_check)
    if len(templateFileName) == 0:
        return False
    elif len(templateFileName) > 1:
        print("More than one template matched. Please check the prefix used or the template names")
    else:
        return templateFileName


def initialize_config_files(path_to_data, path_to_templates):
    filenames = find_data_filenames(path_to_data)
    for filename in filenames:
        if os.path.splitext(filename)[1] == '.yaml':
            continue

        # extractForTemplateCheck gets the prefix of the base filename (without the path). 
        # This is to be used to find the relevant tempalte in templates
        extractForTemplateCheck = os.path.basename(filename).split('_')[0]

        # the configFilename includes the path, the same as filename does
        configFileName = os.path.splitext(filename)[0]+".yaml"

        if configFileName in filenames:
            print(filename + " already has the corresponding config file: " + configFileName)
        else:
            template = check_if_template(extractForTemplateCheck, path_to_templates)
            if template == False:
                configFile = open(configFileName, "w")
                configFile.close()
                print("Created " + configFileName + " as empty config file for " + filename)
            else:
                shutil.copy(template[0], configFileName)
                print("Created " + configFileName + " config file from " + template[0] + " template, for " + filename)
    print("Finished config files initialization")


def fileToPandas(fileName, fileExtension, skipRows=0):
    if fileExtension == '.csv':
        return pd.read_csv(fileName, skiprows=skipRows)
    else:
        return pd.read_excel(fileName, skipRows=skipRows)


dataFolder = "./data"
try:
    pmFolder = argv[1]
    pmDataPath = dataFolder + "/" + pmFolder
except IndexError as e:
    print("PM Data folder not provided as argument")
    pmFolder = ""
    pmDataPath = ""
while True:
    if os.path.isdir(pmDataPath):
        print("Folder containing PM Data: " + pmDataPath)
        break
    else:
        print("PM Data folder " + pmFolder + " doesn't exist in " + dataFolder)
        pmFolder = input('Please enter the PM Data folder name located in ' + dataFolder + ":")
        pmDataPath = dataFolder + "/" + pmFolder
templatesPath = "./templates"
filesExtension = '.csv'
initialize_config_files(pmDataPath, templatesPath)
fileNames = find_data_filenames(pmDataPath, suffix=filesExtension)
skippedRowsNumber = 7

# prepare presentation helpers
prs = Presentation()
outputPath = './output/'
tmpPicturePath = './output/tmp/'
pptName = 'Report_' + mdates.datetime.date.today().strftime('%d-%m-%Y') + '.pptx'
slidesOrder = {}

# prepare figure helpers
panels = []
output_file('Report_' + mdates.datetime.date.today().strftime('%d-%m-%Y') + '.html',
            title='Report_' + mdates.datetime.date.today().strftime('%d-%m-%Y') )

for fileName in fileNames:
    print(fileName)
    configFileName = os.path.splitext(fileName)[0] + ".yaml"

    # check if the config file is empty. If, so skip this excel file
    if os.path.getsize(configFileName) != 0:
        configFile = open(configFileName, "r")
        configDictsList = yaml.load(configFile)
    else:
        print(fileName + " has an empty config file")
        continue
    # check if 'Result' is in the first line of the file. If so, remove 7 lines

    while len(configDictsList) > 0:
        line = open(fileName).readline()
        if 'Result' in line:
            excelfile = fileToPandas(fileName, filesExtension, skippedRowsNumber)
        else:
            excelfile = fileToPandas(fileName, filesExtension)
        # Since the yaml config file contains a list of dictionaries (each file can be used for multiple graphs),
        # we pop each config and create the graph

        configDict = configDictsList.pop(0)
        # change Reported Time type from objet to datetime and set it as index
        print("configDict este: ")
        print(configDict)
        excelfile[configDict['indexul']] = pd.to_datetime(excelfile[configDict['indexul']])
        excelfile.set_index(configDict['indexul'], inplace=True)
        FilterParams = configDict['FilterParams']
        TargetColumns = configDict['TargetColumns']
        UseTargetColumnsInLegend = configDict['UseTargetColumnsInLegend']
        reverseNeeded = configDict['reverseNeeded']

        if reverseNeeded:
            # reverse the oreder from oldest to newest

            excelfile = excelfile[::-1]
        figura = figure(x_axis_type='datetime', plot_width=900)
        legenda = Legend()
        colectieGraphs = GraficCollection(figura)
        colectieGraphs.extrage_gafice(excelfile, FilterParams, TargetColumns, UseTargetColumnsInLegend)
        i=0
        for label, serie in colectieGraphs.grafice.items():
            x = serie.index
            y = serie.values
            graficCurent = colectieGraphs.figura.line(x=x, y=y, color=palletes.Paired[len(colectieGraphs.grafice)][i])
            legenda.items.extend([LegendItem(label=label, renderers=[graficCurent])])
            i += 1
        colectieGraphs.figura.add_layout(legenda, 'above')
        panels.append(Panel(child=colectieGraphs.figura, title=TargetColumns.__str__()))

        # adjustments to the figure
        # ---------

        # colectieGraphs.figura.axes[0].set_xlim(excelfile.index[0], excelfile.index[-1])
        # dateTimeFmt = mdates.DateFormatter('%D %H:%M')
        # colectieGraphs.figura.axes[0].xaxis.set_major_locator(plt.MaxNLocator(45))
        # colectieGraphs.figura.axes[0].xaxis.set_major_formatter(dateTimeFmt)
        # colectieGraphs.figura.axes[0].xaxis.set_tick_params(rotation=90)
        #colectieGraphs.figura.legend(loc=9, mode='none', fontsize='small', ncol=5, labelspacing=0.05)
        # colectieGraphs.figura.legend(loc=9, ncol=5)
        # colectieGraphs.figura.set_size_inches(9.99, 6.7)
        # colectieGraphs.figura.tight_layout()
        # colectieGraphs.figura.subplots_adjust(top=0.900, bottom=0.3)
        # colectieGraphs.figura.add_axes([0.1, 0.01, 0.3, 0.07])
        # axaTextBox = colectieGraphs.figura.axes[-1]
        # text_box = TextBox(axaTextBox, 'Evaluate', initial=FilterParams.__repr__())

        # ----------

        tmpPictureName = configDict['Title'] + '.png'
#        export_png(colectieGraphs.figura, tmpPicturePath + tmpPictureName)

        # Create a dictionary where the key is the slide number and the value is the picture name
        # But first check if the order configured is correct. If not, terminate the script.
        # note that the tmptPictureName will also be used for the Title of the slide
        if configDict['slideNumber'] in slidesOrder.keys():
            print("the slide number already exists. Please check the configs. Current config file: " + configFileName)
            exit()
        slidesOrder[configDict['slideNumber']] = tmpPictureName

tabs = Tabs(tabs=panels)
show(tabs)

# We sort the slidesOrder so the outputed ppt would be in the requred
#sortedKeysAsInt = sorted([int(k) for k in slidesOrder.keys()])
#for slideNumber in sortedKeysAsInt:
    # create a slide and add it to the ppt
#    title_slide_layout = prs.slide_layouts[5]
#    slide = prs.slides.add_slide(title_slide_layout)
#    titlu = slide.placeholders[0]
    # As the value stored for the key is the tmpPictureName, we use it for the title by removing the extension
#    titlu.text = os.path.splitext(slidesOrder[str(slideNumber)])[0]
#    titlu.top = 0
#    titlu.left = 0
#    titlu.height = Inches(0.5)
#    titlu.width = Inches(10)
#    titlu.text_frame.paragraphs[0].font.size = Pt(14)
#    titlu.text_frame.paragraphs[0].font.bold = True
#    pic = slide.shapes.add_picture(tmpPicturePath + slidesOrder[str(slideNumber)], Inches(0.01), Inches(0.5))



#prs.save(outputPath + pptName)



#plt.close()